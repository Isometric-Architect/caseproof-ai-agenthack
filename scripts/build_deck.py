"""Build a compact PowerPoint deck for Devpost."""

from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"


def add_bullets(slide, bullets):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(8.7), Inches(4.6))
    frame = box.text_frame
    frame.clear()
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(23)
        p.space_after = Pt(9)


def add_title(slide, title, subtitle=None):
    from pptx.util import Inches, Pt

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.05), Inches(8.7), Inches(0.4))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(17)


def main() -> int:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover
        raise SystemExit(f"python-pptx is required: {exc}")

    DOCS.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = [
        (
            "CaseProof AI",
            "A human-review gate for high-value refund exception cases.",
            [
                "Agents can prepare refund cases, but preparation is not approval.",
                "CaseProof checks evidence, policy, required milestones, and review routing.",
                "The only allowed outcome is a human review task.",
            ],
        ),
        (
            "Problem",
            None,
            [
                "A refund case can look complete while delivery proof or a customer statement is missing.",
                "An agent can skip policy or risk review.",
                "A high-value refund can be routed as if it were low risk.",
            ],
        ),
        (
            "What It Checks",
            None,
            [
                "Required evidence is present and referenced.",
                "The decision belongs to the same case.",
                "Required Maestro case milestones are present.",
                "Policy caps route risky cases to human review.",
            ],
        ),
        (
            "Demo Outcomes",
            None,
            [
                "HOLD: refund evidence is missing.",
                "ALLOW_HUMAN_REVIEW_ONLY: ready for a person.",
                "BLOCK: policy bypass or skipped milestone.",
                "Auto-approval remains disabled.",
            ],
        ),
        (
            "UiPath Fit",
            None,
            [
                "UiPath Maestro is the orchestration layer.",
                "Studio Web or API Workflow can call the validator.",
                "The result opens a review task instead of taking an external action.",
                "Public repo uses synthetic packets only.",
            ],
        ),
        (
            "Claim Boundary",
            None,
            [
                "This is a public-safe hackathon prototype.",
                "It is not production approval or compliance certification.",
                "A real tenant run should bind case URL, process instance, and review task.",
            ],
        ),
    ]

    for title, subtitle, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets)

    if (ASSETS / "caseproof_dashboard_desktop.png").exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Demo Dashboard")
        slide.shapes.add_picture(str(ASSETS / "caseproof_dashboard_desktop.png"), Inches(0.75), Inches(1.15), width=Inches(8.55))

    prs.save(DOCS / "CaseProof_AI_AgentHack_Deck_Draft.pptx")
    print("CASEPROOF_DECK_BUILT")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
